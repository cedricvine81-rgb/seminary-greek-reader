# -*- coding: utf-8 -*-
"""
Apply the agreed corrections to the lesson PowerPoints.

    python3 scripts/fix-decks.mjs.py ~/dev/decks/"Lessons (seminarygreek)"

Each edit is a (file, slide, old, new) quadruple, applied by replacing the span ACROSS the
runs it covers rather than rewriting the paragraph — PowerPoint has split these lines into a
run per accented character, and collapsing them would take the formatting with it. Runs that
do not overlap the span are not touched at all.

Scope, as agreed with the instructor:
  - the seven slips where the slide's Greek is not a Greek word or is the wrong form
  - accents added to the quotations on those slides (the decks are otherwise unaccented)
  - Rom 14:19 replaced: our text prints the INDICATIVE διώκομεν, so it cannot illustrate the
    hortatory subjunctive the slide is teaching
  - the four edition differences deliberately LEFT ALONE (εἶδον, Πιλᾶτος, ἦν, ἔπαθεν)
"""
import sys, os
from pptx import Presentation

EDITS = [
  # ── Lesson 7 · Using Participles ──
  ("Lesson 7/2. Lesson 7 (Using Participles).pptx", 32,
   "ἀκουσασα περι του Ἰησου, ἐλθουσα ἐν τῳ ὀχλῳ ὀπισθεν ἡψατα του ἱματιου αὐτου.",
   "ἀκούσασα περὶ τοῦ Ἰησοῦ, ἐλθοῦσα ἐν τῷ ὄχλῳ ὄπισθεν ἥψατο τοῦ ἱματίου αὐτοῦ."),
  ("Lesson 7/2. Lesson 7 (Using Participles).pptx", 49,
   "οὐκ εἰμι ἱκανος κυφας λυσαι τον",
   "οὐκ εἰμὶ ἱκανὸς κύψας λῦσαι τὸν"),
  ("Lesson 7/2. Lesson 7 (Using Participles).pptx", 49,
   "ἱμαντα των ὑποδηματων αὐτου.",
   "ἱμάντα τῶν ὑποδημάτων αὐτοῦ."),

  # ── Lesson 8 · Subjunctive ──  (ἔπαθεν kept: an edition difference, not a slip)
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 36,
   "Χριστος ἁπαξ περι ἁματιων ἐπαθεν",
   "Χριστὸς ἅπαξ περὶ ἁμαρτιῶν ἔπαθεν"),
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 36,
   "ἱνα ὑμας",
   "ἵνα ὑμᾶς"),
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 36,
   "προσαγαγῃ τῳ θεῳ.",
   "προσαγάγῃ τῷ θεῷ."),
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 36,
   "ὁς ἀν ἐσθιῃ τον ἀρτον ἠ πινῃ το ποτηριον του",
   "ὃς ἂν ἐσθίῃ τὸν ἄρτον ἢ πίνῃ τὸ ποτήριον τοῦ"),
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 36,
   "κυριος ...",
   "κυρίου ..."),
  # Rom 14:19 prints the indicative in our text — swapped for an unambiguous hortatory
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 28,
   "Rom. 14.19 ἀρα οὐν τα της εἰρηνης διωκωμεν.",
   "1 John 4.7: ἀγαπῶμεν ἀλλήλους."),
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 28,
   "– So therefore let us pursue the things of peace.",
   "– Let us love one another."),
  # accents on the two hortatory examples that stay
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 28,
   "διελθωμεν εἰς το περαν.", "διέλθωμεν εἰς τὸ πέραν."),
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 28,
   "προσερχωμεθα μετα ἀληθινης καρδιας ἐν πληροφοριᾳ",
   "προσερχώμεθα μετὰ ἀληθινῆς καρδίας ἐν πληροφορίᾳ"),
  ("Lesson 8/1. Lesson 8 (Subjunctive).pptx", 28, "πιστεως.", "πίστεως."),

  # ── Lesson 9 · Mi Verbs ──
  ("Lesson 9/Lesson 9 (Mi Verbs).pptx", 37,
   "ἐλεγεν αὐτοις∙ Ὑμιν το μυστηριον δεδοται της βασιλειας",
   "ἔλεγεν αὐτοῖς· ὑμῖν τὸ μυστήριον δέδοται τῆς βασιλείας"),
  ("Lesson 9/Lesson 9 (Mi Verbs).pptx", 37, "του θεον.", "τοῦ θεοῦ."),
  ("Lesson 9/Lesson 9 (Mi Verbs).pptx", 38,
   "και στας ὁ Ἰησους ἐφωνησαν αὐτους και εἰπεν",
   "καὶ στὰς ὁ Ἰησοῦς ἐφώνησεν αὐτοὺς καὶ εἶπεν"),
  ("Lesson 9/Lesson 9 (Mi Verbs).pptx", 46,
   "και ἐμετιμησεν αὐτοις ἱνα μη φανερον αὐτον ποιησωσιν, ἱνα πληρωθῃ το ῥηθεν δια Ἠσαϊου του προφητου",
   "καὶ ἐπετίμησεν αὐτοῖς ἵνα μὴ φανερὸν αὐτὸν ποιήσωσιν, ἵνα πληρωθῇ τὸ ῥηθὲν διὰ Ἡσαΐου τοῦ προφήτου"),
]


def replace_span(para, old, new):
    """Replace `old` with `new` across whatever runs it spans, leaving the others untouched."""
    runs = para.runs
    full = ''.join(r.text for r in runs)
    i = full.find(old)
    if i < 0:
        return False
    j = i + len(old)
    pos, placed = 0, False
    for r in runs:
        s, e = pos, pos + len(r.text)
        pos = e
        if e <= i or s >= j:          # wholly outside the span
            continue
        before = r.text[:i - s] if s < i else ''
        after = r.text[j - s:] if e > j else ''
        if not placed:
            r.text = before + new + after
            placed = True
        else:
            r.text = before + after
    return placed


def main(root):
    by_file = {}
    for rel, slide, old, new in EDITS:
        by_file.setdefault(rel, []).append((slide, old, new))

    for rel, edits in by_file.items():
        path = os.path.join(root, rel)
        prs = Presentation(path)
        done = failed = 0
        for slide_no, old, new in edits:
            hit = False
            for shape in prs.slides[slide_no - 1].shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    if replace_span(para, old, new):
                        hit = True
                        break
                if hit:
                    break
            if hit:
                done += 1
            else:
                failed += 1
                print(f"  MISS  {rel} s{slide_no}: {old[:60]}")
        prs.save(path)
        print(f"  {rel}: {done} applied{f', {failed} MISSED' if failed else ''}")


if __name__ == '__main__':
    main(sys.argv[1] if len(sys.argv) > 1 else os.path.expanduser('~/dev/decks/Lessons (seminarygreek)'))
